import math
import os
import sys
import time

from PyQt4 import QtGui, QtCore

from matplotlib.backends.backend_qt4agg import FigureCanvasQTAgg, NavigationToolbar2QT
from matplotlib.ticker import ScalarFormatter

import matplotlib as mpl
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches

from config import *
from dat_file import DatFile
from operations import Operations

"""
- Axis scaling filters
- Low pass filter
- Integration into qtlab as real time plotting
- Look into plotting performance
- Use proper axis when plotting linecut and not averaging
"""

class FixedOrderFormatter(ScalarFormatter):
    def __init__(self, significance=0):
        ScalarFormatter.__init__(self, useOffset=None, useMathText=None)
        self.format = '%.' + str(significance) + 'f'

    def __call__(self, x, pos=None):
        if x == 0:
            return '0'

        exp = self.orderOfMagnitude

        return self.format % (x / (10 ** exp))

    def _set_orderOfMagnitude(self, range):
        exp = math.floor(math.log10(range))
        self.orderOfMagnitude = exp - (exp % 3)

class Window(QtGui.QMainWindow):
    def __init__(self, lc_window, op_window, filename=None):
        QtGui.QMainWindow.__init__(self)
        
        self.linecut = lc_window
        self.operations = op_window

        self.fig, self.ax = plt.subplots()
        self.cb = None

        self.line = None
        self.linecut_type = None
        self.linecut_coord = None

        self.ppt = None
        self.slide = None

        self.data = None
        self.data_file = None

        self.axis_changed = False

        self.init_ui()

        if filename is not None:
            self.load_file(filename)
            self.update_ui()

    def init_ui(self):
        self.setWindowTitle('qtplot')

        self.main_widget = QtGui.QWidget(self)

        self.canvas = FigureCanvasQTAgg(self.fig)
        self.canvas.mpl_connect('button_press_event', self.on_mouse_click)
        self.canvas.mpl_connect('motion_notify_event', self.on_mouse_motion)
        self.toolbar = NavigationToolbar2QT(self.canvas, self)

        self.b_load = QtGui.QPushButton('Load DAT...')
        self.b_load.clicked.connect(self.on_load_dat)
        self.b_swap = QtGui.QPushButton('Swap order', self)
        self.b_swap.clicked.connect(self.on_swap_order)
        self.c_average = QtGui.QCheckBox('Average Y-Axis', self)
        self.c_average.setChecked(True)
        self.c_average.stateChanged.connect(self.on_data_change)

        self.lbl_x = QtGui.QLabel("X", self)
        self.cb_x = QtGui.QComboBox(self)
        self.cb_x.activated.connect(self.on_axis_changed)
        lbl_order_x = QtGui.QLabel('Order: ', self)
        self.cb_order_x = QtGui.QComboBox(self)

        self.lbl_y = QtGui.QLabel("Y", self)
        self.cb_y = QtGui.QComboBox(self)
        self.cb_y.activated.connect(self.on_axis_changed)
        lbl_order_y = QtGui.QLabel('Order: ', self)
        self.cb_order_y = QtGui.QComboBox(self)

        self.lbl_d = QtGui.QLabel("Data", self)
        self.cb_z = QtGui.QComboBox(self)
        self.cb_z.activated.connect(self.on_axis_changed)

        self.le_min = QtGui.QLineEdit(self)
        self.le_min.returnPressed.connect(self.on_cmap_changed)

        self.s_gamma = QtGui.QSlider(QtCore.Qt.Horizontal)
        self.s_gamma.setMinimum(0)
        self.s_gamma.setMaximum(100)
        self.s_gamma.setValue(50)
        self.s_gamma.valueChanged.connect(self.on_cmap_changed)

        self.le_max = QtGui.QLineEdit(self)
        self.le_max.returnPressed.connect(self.on_cmap_changed)

        self.lbl_ppt = QtGui.QLabel("PPT File", self)
        self.b_ppt = QtGui.QPushButton("Browse...", self)
        self.b_ppt.clicked.connect(self.browse_ppt)
        self.le_ppt = QtGui.QLineEdit(cfg_ppt_file, self)

        self.b_slide = QtGui.QPushButton('Add slide')
        self.b_slide.setEnabled(False)
        self.b_slide.clicked.connect(self.add_slide)

        self.b_ppt1 = QtGui.QPushButton('Add 2D data Figure')
        self.b_ppt1.setEnabled(False)
        self.b_ppt1.clicked.connect(self.add_2d_data)

        self.b_ppt2 = QtGui.QPushButton('Add linecut figure')
        self.b_ppt2.setEnabled(False)
        self.b_ppt2.clicked.connect(self.add_linecut)

        self.b_refresh = QtGui.QPushButton('Save PPT')
        self.b_refresh.clicked.connect(self.save_ppt)

        hbox = QtGui.QHBoxLayout()
        hbox.addWidget(self.b_load)
        hbox.addWidget(self.b_swap)
        hbox.addWidget(self.c_average)
        
        hbox1 = QtGui.QHBoxLayout()
        hbox1.addWidget(self.lbl_x)
        hbox1.addWidget(self.cb_x)
        hbox1.addWidget(lbl_order_x)
        hbox1.addWidget(self.cb_order_x)

        hbox2 = QtGui.QHBoxLayout()
        hbox2.addWidget(self.lbl_y)
        hbox2.addWidget(self.cb_y)
        hbox2.addWidget(lbl_order_y)
        hbox2.addWidget(self.cb_order_y)

        hbox3 = QtGui.QHBoxLayout()
        hbox3.addWidget(self.lbl_d)
        hbox3.addWidget(self.cb_z)

        hbox_gamma = QtGui.QHBoxLayout()
        hbox_gamma.addWidget(self.le_min)
        hbox_gamma.addWidget(self.s_gamma)
        hbox_gamma.addWidget(self.le_max)

        hbox4 = QtGui.QHBoxLayout()
        hbox4.addWidget(self.lbl_ppt)
        hbox4.addWidget(self.b_ppt)
        hbox4.addWidget(self.le_ppt)

        hbox5 = QtGui.QHBoxLayout()
        hbox5.addWidget(self.b_slide)
        hbox5.addWidget(self.b_ppt1)
        hbox5.addWidget(self.b_ppt2)
        hbox5.addWidget(self.b_refresh)

        vbox = QtGui.QVBoxLayout(self.main_widget)
        vbox.addWidget(self.toolbar)
        vbox.addWidget(self.canvas)
        vbox.addLayout(hbox)
        vbox.addLayout(hbox1)
        vbox.addLayout(hbox2)
        vbox.addLayout(hbox3)
        vbox.addLayout(hbox_gamma)
        vbox.addLayout(hbox4)
        vbox.addLayout(hbox5)

        self.main_widget.setFocus()
        self.setCentralWidget(self.main_widget)

        self.move(100, 100)

    def update_ui(self):
        self.setWindowTitle(self.name)

        self.cb_x.clear()
        self.cb_x.addItems(self.data_file.columns)
        self.cb_x.setCurrentIndex(cfg_default_x)
        self.cb_order_x.addItems(self.data_file.columns)
        self.cb_order_x.setCurrentIndex(0)

        self.cb_y.clear()
        self.cb_y.addItems(self.data_file.columns)
        self.cb_y.setCurrentIndex(cfg_default_y)
        self.cb_order_y.addItems(self.data_file.columns)
        self.cb_order_y.setCurrentIndex(1)

        self.cb_z.clear()
        self.cb_z.addItems(self.data_file.columns)
        self.cb_z.setCurrentIndex(cfg_default_data)

        self.b_slide.setEnabled(True)
        self.b_ppt1.setEnabled(True)
        self.b_ppt2.setEnabled(True)

    def load_file(self, filename):
        if filename != "":
            self.data_file = DatFile(filename)

            path, self.name = os.path.split(self.data_file.filename)

            self.update_ui()
            self.on_data_change()

    def on_load_dat(self, event):
        self.filename = str(QtGui.QFileDialog.getOpenFileName(filter='*.dat'))

        self.load_file(self.filename)

    def on_swap_order(self, event):
        x, y = self.cb_order_x.currentIndex(), self.cb_order_y.currentIndex()
        self.cb_order_x.setCurrentIndex(y)
        self.cb_order_y.setCurrentIndex(x)

        self.on_data_change()

    def on_axis_changed(self, event):
        self.axis_changed = True
        self.on_data_change()

    def on_cmap_changed(self):
        self.plot_2d_data()

    def on_mouse_click(self, event):
        if not event.inaxes:
            return

        self.plot_linecut(event.button, event.xdata, event.ydata)

    def on_mouse_motion(self, event):
        if event.button != None:
            self.on_mouse_click(event)

    def on_data_change(self):
        if self.data_file is not None:
            self.manipulate_data()
            self.plot_2d_data()

    # Get a matrix of the data values with x and y axes
    def get_averaged_pivot(self, data, x, y, z, order_x, order_y):
        processed = data.copy()

        if x != order_x:
            processed[x] = processed.groupby(order_x)[x].transform(np.average)

        if y != order_y:
            processed[y] = processed.groupby(order_y)[y].transform(np.average)

        return processed.pivot(index=y, columns=x, values=z)

    def manipulate_data(self):
        # Get the column names to use for the x, y and data
        self.lbl_x = str(self.cb_x.currentText())
        self.lbl_y = str(self.cb_y.currentText())
        self.data_lbl = str(self.cb_z.currentText())

        order_x, order_y = str(self.cb_order_x.currentText()), str(self.cb_order_y.currentText())

        data = self.get_averaged_pivot(self.data_file.df, self.lbl_x, self.lbl_y, self.data_lbl, order_x, order_y)

        self.data = self.operations.apply_operations(data)

        if not self.c_average.isChecked():
            self.y_coords = self.data_file.df.pivot(index=order_y, columns=order_x, values=self.lbl_y)

        self.data_changed = False

    # As long as axis coords and values are in the same order
    def get_quadrilaterals(self, data, order_x, order_y, axis_y, average_y):
        xc, yc = np.meshgrid(data.columns.values, data.index.values)

        # TODO remove dependency on state
        if not average_y:
            yc = self.data_file.df.pivot(index=order_y, columns=order_x, values=axis_y).values

        # Add a column/row on each side with estimated values
        xc = np.hstack((xc[:,[0]] - (xc[:,[1]] - xc[:,[0]]), xc, xc[:,[-1]] + (xc[:,[-1]] - xc[:,[-2]])))
        yc = np.vstack([yc[0] - (yc[1] - yc[0]), yc, yc[-1] + (yc[-1] - yc[-2])])

        x = xc[:,:-1] + np.diff(xc, axis=1) / 2.0
        y = yc[:-1,:] + np.diff(yc, axis=0) / 2.0

        x = np.vstack((x, x[-1]))
        y = np.hstack([y, y[:,[-1]]])

        return x, y

    def plot_2d_data(self):
        if self.data is None:
            return

        cmap = mpl.cm.get_cmap('seismic')
        value = (self.s_gamma.value() / 100.0)
        cmap.set_gamma(math.exp(value * 5) / 10.0)

        self.ax.clear()

        order_x, order_y = str(self.cb_order_x.currentText()), str(self.cb_order_y.currentText())
        x, y = self.get_quadrilaterals(self.data, order_x, order_y, self.lbl_y, self.c_average.isChecked())

        # Mask NaN values so they will not be plotted
        masked_y = np.ma.masked_where(np.isnan(y), y)
        masked = np.ma.masked_where(np.isnan(self.data.values), self.data.values)
        self.quadmesh = self.ax.pcolormesh(x, masked_y, masked, cmap=cmap)

        if self.le_min.text() == '' or self.le_max.text() == '' or self.axis_changed:
            cm_min, cm_max = self.quadmesh.get_clim()
            self.le_min.setText('%.2e' % cm_min)
            self.le_max.setText('%.2e' % cm_max)

            self.s_gamma.setValue(50)
        else:
            self.quadmesh.set_clim(vmin=float(self.le_min.text()), vmax=float(self.le_max.text()))
        
        self.ax.axis('tight')

        # Create a colorbar, if there is already one draw it in the existing place
        if self.cb:
            self.cb.ax.clear()
            self.cb = self.fig.colorbar(self.quadmesh, cax=self.cb.ax)
        else:
            self.cb = self.fig.colorbar(self.quadmesh)

        self.cb.set_label(self.data_lbl)
        self.cb.formatter = FixedOrderFormatter(1)
        self.cb.update_ticks()

        # Set the various labels
        self.ax.set_title(self.name)
        self.ax.set_xlabel(self.lbl_x)
        self.ax.set_ylabel(self.lbl_y)
        self.ax.xaxis.set_major_formatter(FixedOrderFormatter())
        self.ax.yaxis.set_major_formatter(FixedOrderFormatter())
        self.ax.set_aspect('auto')

        self.fig.tight_layout()
        
        self.canvas.draw()

        self.background = self.canvas.copy_from_bbox(self.ax.bbox)

        self.axis_changed = False

    def plot_linecut(self, button, x_coord, y_coord):
        lc = self.linecut

        if len(self.ax.lines) == 0:
            self.line = self.ax.axvline(color='red')

        if len(lc.ax.lines) > 0:
            lc.ax.lines.pop(0)

        if button == 1:
            # Get the row closest to the mouse Y
            self.linecut_coord = min(self.data.index, key=lambda x:abs(x - y_coord))
            self.linecut_type = 'horizontal'

            self.line.set_transform(self.ax.get_yaxis_transform())
            self.line.set_xdata([0, 1])
            self.line.set_ydata([self.linecut_coord, self.linecut_coord])

            lc.ax.plot(self.data.columns, self.data.loc[self.linecut_coord], color='red', linewidth=0.5)
            lc.ax.set_xlabel(self.lbl_x)
        elif button == 2:
            # Get the column closest to the mouse X
            self.linecut_coord = min(self.data.columns, key=lambda x:abs(x - x_coord))
            self.linecut_type = 'vertical'

            self.line.set_transform(self.ax.get_xaxis_transform())
            self.line.set_xdata([self.linecut_coord, self.linecut_coord])
            self.line.set_ydata([0, 1])

            lc.ax.plot(self.data.index, self.data[self.linecut_coord], color='red', linewidth=0.5)
            lc.ax.set_xlabel(self.lbl_y)
        
        lc.ax.set_title(self.name)
        lc.ax.set_ylabel(self.data_lbl)
        lc.ax.xaxis.set_major_formatter(FixedOrderFormatter())
        lc.ax.yaxis.set_major_formatter(FixedOrderFormatter())

        # Make the figure fit the data
        lc.ax.relim()
        lc.ax.autoscale_view()
        lc.fig.tight_layout()

        # Redraw both plots to update them
        self.canvas.restore_region(self.background)
        self.ax.draw_artist(self.line)
        self.canvas.blit(self.ax.bbox)

        lc.canvas.draw()

    def add_slide(self, event):
        if self.ppt is not None:
            title_slide_layout = self.ppt.slide_layouts[6]
            self.slide = self.ppt.slides.add_slide(title_slide_layout)

    def add_2d_data(self, event):
        if self.slide is not None:
            path = os.path.dirname(os.path.realpath(__file__))
            path = os.path.join(path, 'test.png')
            self.fig.savefig(path)
            self.slide.shapes.add_picture(path, Inches(1), Inches(1))

    def add_linecut(self, event):
        if self.slide is not None:
            path = os.path.dirname(os.path.realpath(__file__))
            path = os.path.join(path, 'test.png')
            self.linecut.fig.savefig(path)
            self.slide.shapes.add_picture(path, Inches(1), Inches(1))

    def browse_ppt(self, event):
        filename = str(QtGui.QFileDialog.getOpenFileName(self))
        self.le_ppt.setText(filename)

        try:
            self.ppt_file = open(filename, 'rb')
            self.ppt = Presentation(self.ppt_file)
        except IOError:
            print 'Could not open PowerPoint file: ' + filename
            self.ppt = None

    def save_ppt(self, event):
        if self.ppt:
            try:
                self.ppt.save(str(self.le_ppt.text()))
            except IOError as e:
                print 'Could not save PowerPoint file: '

    def resizeEvent(self, event):
        if len(self.ax.lines) > 0:
            self.ax.lines.pop(0)

        # Very slow, maybe search for faster way
        # http://stackoverflow.com/questions/13552345/how-to-disable-multiple-auto-redrawing-at-resizing-widgets-in-pyqt
        self.plot_2d_data()

    def closeEvent(self, event):
        if self.ppt:
            try:
                self.ppt.save(str(self.le_ppt.text()))
            except IOError as e:
                path, filename = os.path.split(str(self.le_ppt.text()))
                name, ext = os.path.splitext(filename)
                new = os.path.join(path, (name + time.asctime()) + ext)
                print new
                self.ppt.save(new)

        self.linecut.close()
        self.operations.close()

class Linecut(QtGui.QDialog):
    def __init__(self, parent=None):
        super(Linecut, self).__init__(parent)

        self.fig, self.ax = plt.subplots()

        self.init_ui()

    def init_ui(self):
        self.setWindowTitle("Linecut")

        self.canvas = FigureCanvasQTAgg(self.fig)
        self.toolbar = NavigationToolbar2QT(self.canvas, self)

        layout = QtGui.QVBoxLayout()
        layout.addWidget(self.toolbar)
        layout.addWidget(self.canvas)
        self.setLayout(layout)

        self.move(800, 100)

if __name__ == '__main__':
    app = QtGui.QApplication(sys.argv)

    linecut = Linecut()
    operations = Operations()
    main = Window(linecut, operations)
    
    if len(sys.argv) > 1:
        main = Window(linecut, operations, filename=sys.argv[1])

    linecut.main = main
    operations.main = main

    linecut.show()
    main.show()

    operations.show()

    sys.exit(app.exec_())